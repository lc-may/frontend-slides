#!/usr/bin/env python3
"""
PowerPoint Content Extractor for Frontend Slides

Extracts all content from PPT/PPTX files for conversion to HTML presentations.
Outputs JSON structure with slides, text, images, and speaker notes.

Usage:
    python pptx-extract.py <input.pptx> <output_dir>

Example:
    python pptx-extract.py presentation.pptx ./extracted
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import json
import os
import sys


def extract_pptx(file_path, output_dir):
    """
    Extract all content from a PowerPoint file.

    Returns a JSON structure with slides, text, and images.

    Args:
        file_path: Path to the .pptx file
        output_dir: Directory to save extracted assets

    Returns:
        List of slide data dictionaries
    """
    prs = Presentation(file_path)
    slides_data = []

    # Create assets directory
    assets_dir = os.path.join(output_dir, 'assets')
    os.makedirs(assets_dir, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides):
        slide_data = {
            'number': slide_num + 1,
            'title': '',
            'content': [],
            'images': [],
            'notes': ''
        }

        for shape in slide.shapes:
            # Extract title
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    slide_data['title'] = shape.text
                else:
                    slide_data['content'].append({
                        'type': 'text',
                        'content': shape.text
                    })

            # Extract images
            if shape.shape_type == 13:  # Picture
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                image_name = f"slide{slide_num + 1}_img{len(slide_data['images']) + 1}.{image_ext}"
                image_path = os.path.join(assets_dir, image_name)

                with open(image_path, 'wb') as f:
                    f.write(image_bytes)

                slide_data['images'].append({
                    'path': f"assets/{image_name}",
                    'width': shape.width,
                    'height': shape.height
                })

        # Extract notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            slide_data['notes'] = notes_frame.text

        slides_data.append(slide_data)

    return slides_data


def main():
    if len(sys.argv) < 3:
        print("Usage: python pptx-extract.py <input.pptx> <output_dir>")
        print("Example: python pptx-extract.py presentation.pptx ./extracted")
        sys.exit(1)

    input_file = sys.argv[1]
    output_dir = sys.argv[2]

    if not os.path.exists(input_file):
        print(f"Error: File not found: {input_file}")
        sys.exit(1)

    # Create output directory
    os.makedirs(output_dir, exist_ok=True)

    # Extract content
    print(f"Extracting content from {input_file}...")
    slides_data = extract_pptx(input_file, output_dir)

    # Save JSON output
    output_json = os.path.join(output_dir, 'slides.json')
    with open(output_json, 'w', encoding='utf-8') as f:
        json.dump(slides_data, f, indent=2, ensure_ascii=False)

    # Print summary
    print(f"\nExtraction complete!")
    print(f"  Slides: {len(slides_data)}")
    print(f"  Output: {output_dir}/")

    total_images = sum(len(s['images']) for s in slides_data)
    if total_images > 0:
        print(f"  Images: {total_images} (saved to {output_dir}/assets/)")

    print(f"\nSlide breakdown:")
    for slide in slides_data:
        img_count = len(slide['images'])
        title = slide['title'] or "(no title)"
        print(f"  Slide {slide['number']}: {title}")
        if img_count:
            print(f"    └─ {img_count} image(s)")

    return slides_data


if __name__ == "__main__":
    main()
